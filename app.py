import streamlit as st
import base64
import groq
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from io import BytesIO
from dotenv import load_dotenv

load_dotenv()

# Configure Groq API
genai = groq.Client(api_key=os.getenv('GROQ_API_KEY'))

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)
MAX_CONTENT_LENGTH = 300

# Define margin size for the content box (in inches)
CONTENT_LEFT = Inches(1)
CONTENT_TOP = Inches(2)
CONTENT_WIDTH = Inches(8.5)
CONTENT_HEIGHT = Inches(5)

def generate_slide_titles(topic):
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = genai.chat.completions.create(
        model="llama3-70b-8192",
        messages=[{"role": "system", "content": "You are an expert in creating slide presentations."},
                  {"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response.choices[0].message.content.split("\n")

def generate_slide_content(slide_title):
    prompt = f"Generate concise and structured content for the slide: '{slide_title}'."
    response = genai.chat.completions.create(
        model="llama3-70b-8192",
        messages=[{"role": "system", "content": "You are an expert in summarizing content for PowerPoint slides."},
                  {"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response.choices[0].message.content[:MAX_CONTENT_LENGTH]

def apply_theme(slide, theme):
    theme_colors = {
        "Light": (RGBColor(242, 242, 242), RGBColor(0, 0, 0), RGBColor(50, 50, 50)),
        "Dark": (RGBColor(0, 0, 0), RGBColor(255, 255, 255), RGBColor(200, 200, 200)),
        "Blue": (RGBColor(0, 102, 204), RGBColor(255, 255, 255), RGBColor(255, 255, 255)),
        "Default": (RGBColor(255, 255, 255), RGBColor(0, 0, 0), RGBColor(50, 50, 50)),
    }
    background_color, title_color, content_color = theme_colors.get(theme, theme_colors["Default"])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background_color

def create_presentation(topic, slide_titles, slide_contents, theme):
    if not os.path.exists('generated_ppt'):
        os.makedirs('generated_ppt')

    prs = pptx.Presentation()
    title_slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    content_slide_layout = prs.slide_layouts[5]  # Title + Content Layout

    # Title Slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = topic
    apply_theme(title_slide, theme)

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(content_slide_layout)
        apply_theme(slide, theme)

        # Set Title Box
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        
        # Set Content Box
        content_box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
        text_frame = content_box.text_frame
        text_frame.text = slide_content
        text_frame.paragraphs[0].font.size = SLIDE_FONT_SIZE
        text_frame.word_wrap = True

    prs.save(f"generated_ppt/{topic}_presentation.pptx")

def main():
    st.title("PowerPoint Presentation Generator with LLaMA-3 📊")
    topic = st.text_input("Enter the topic for your presentation:")
    theme = st.selectbox("Select a theme for your presentation:", ["Light", "Dark", "Blue", "Default"])
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles = [title.strip() for title in slide_titles if title.strip()]
        slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
        create_presentation(topic, filtered_slide_titles, slide_contents, theme)
        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{topic}_presentation.pptx">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()